import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE_TYPE
import config
import os
import math
import json
import logging
import hashlib
import base64
import io
from PIL import Image
import requests
from dotenv import load_dotenv
import openai  # For future use
from typing import List, Dict, Optional, Tuple
import matplotlib.pyplot as plt

log = logging.getLogger(__name__)

def create_presentation(enriched_slides: list, topic: str, style: str = 'dark', slides: int = 6) -> str:
    print("-> Drawing presentation from scratch...")
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    blank_layout = prs.slide_layouts[slides]
    
    # Enhanced theme configuration
    themes = {
        'dark': {
            'font': 'Calibri',
            'text': RGBColor(255, 255, 255),
            'subtext': RGBColor(200, 200, 200),
            'bg': RGBColor(45, 52, 54),
            'accent': RGBColor(52, 152, 219),
            'secondary': RGBColor(44, 62, 80)
        },
        'light': {
            'font': 'Calibri',
            'text': RGBColor(30, 30, 30),
            'subtext': RGBColor(80, 80, 80),
            'bg': RGBColor(248, 249, 250),
            'accent': RGBColor(41, 128, 185),
            'secondary': RGBColor(52, 73, 94)
        }
    }
    theme = themes.get(style, themes['dark'])

    # Add title slide
    slide = prs.slides.add_slide(blank_layout)
    _draw_title_slide(slide, prs, {'slide_title': topic}, theme, topic)

    # Add table of contents
    slide = prs.slides.add_slide(blank_layout)
    _draw_toc_slide(slide, prs, enriched_slides, theme)

    # Add content slides
    for i, slide_data in enumerate(enriched_slides):
        slide = prs.slides.add_slide(blank_layout)
        layout = slide_data.get('layout', 'Photo Layout')
        
        if "Diagram" in layout and slide_data.get('image_path'):
            _draw_diagram_slide(slide, prs, slide_data, theme)
        else:
            _draw_photo_slide(slide, prs, slide_data, theme)

    safe_topic = re.sub(r'[\\/*?:"<>|]', "", topic).replace(" ", "_")
    output_filename = config.PPTConfig.PATHS['output'] / f"{safe_topic}_{style}_presentation.pptx"
    prs.save(output_filename)
    print(f"-> Majestic presentation saved: {output_filename}")
    return str(output_filename)

def _add_image_as_background(slide, prs, image_path, overlay_alpha=0.2):
    """
    Adds an image as a slide background and places a semi-transparent
    black overlay on top to improve text readability.

    This function should be called BEFORE adding any other content to the slide,
    as it places the image and overlay at the bottom of the shape stack.

    Args:
        slide (pptx.slide.Slide): The slide object to modify.
        prs (pptx.presentation.Presentation): The presentation object, used for dimensions.
        image_path (str): The file path to the background image.
        overlay_alpha (float): The opacity of the black overlay (0.0 to 1.0).
                               0.2 means 20% opaque (80% transparent).

    Returns:
        bool: True on success, False on failure.
    """
    from lxml import etree
    
    # Input validation
    if not image_path:
        print("No background image path provided; skipping background.")
        return False

    if not os.path.exists(image_path):
        print(f"Background image file not found: {image_path}")
        return False
    
    # Validate overlay_alpha range
    if not 0.0 <= overlay_alpha <= 1.0:
        print(f"Invalid overlay_alpha value: {overlay_alpha}. Must be between 0.0 and 1.0.")
        return False

    print(f"Adding background from: {image_path}")
    
    try:
        # Get slide dimensions
        slide_width = prs.slide_width
        slide_height = prs.slide_height
        
        # 1. Add background image shape
        pic = slide.shapes.add_picture(
            image_path,
            left=0, top=0,
            width=slide_width, height=slide_height
        )
        
        # 2. Send the picture to the very back
        try:
            shape_elm = pic.element
            shape_tree = slide.shapes._spTree
            shape_tree.remove(shape_elm)
            shape_tree.insert(2, shape_elm)
        except (AttributeError, IndexError) as e:
            print(f"Warning: Could not reorder background image to back: {e}")
        
        # 3. Add readability overlay (scrim) only if overlay_alpha > 0
        if overlay_alpha > 0:
            overlay = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                left=0, top=0,
                width=slide_width, height=slide_height
            )
            
            # 4. Configure the overlay with solid black fill first
            fill = overlay.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(0, 0, 0)
            
            # 5. Add transparency using direct XML manipulation
            # This is the only reliable way since python-pptx doesn't have transparency API
            try:
                # Get the shape's spPr (shape properties) element
                spPr = overlay.element.spPr
                
                # Find or create the solidFill element
                solidFill = spPr.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill')
                
                if solidFill is not None:
                    # Find or create the srgbClr element
                    srgbClr = solidFill.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
                    
                    if srgbClr is not None:
                        # Remove any existing alpha element
                        existing_alpha = srgbClr.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}alpha')
                        if existing_alpha is not None:
                            srgbClr.remove(existing_alpha)
                        
                        # Add alpha element for transparency
                        # Alpha value is in percentage * 1000 (e.g., 20% = 20000)
                        alpha_percentage = int(overlay_alpha * 100 * 1000)
                        alpha_elem = etree.SubElement(
                            srgbClr, 
                            '{http://schemas.openxmlformats.org/drawingml/2006/main}alpha'
                        )
                        alpha_elem.set('val', str(alpha_percentage))
                        
                        print(f"Set overlay transparency using XML: alpha={overlay_alpha} (val={alpha_percentage})")
                    else:
                        print("Warning: Could not find srgbClr element for transparency")
                else:
                    print("Warning: Could not find solidFill element for transparency")
                    
            except Exception as xml_error:
                print(f"XML transparency manipulation failed: {xml_error}")
                # Fallback: try the gradient method
                try:
                    fill.gradient()
                    fill.gradient_stops[0].color.rgb = RGBColor(0, 0, 0)
                    # Note: This might not work either, but it's worth trying
                    print("Used gradient fallback (may not show transparency)")
                except Exception as gradient_error:
                    print(f"Gradient fallback also failed: {gradient_error}")
            
            # 6. Remove the outline from the overlay shape
            try:
                overlay.line.fill.background()
            except AttributeError:
                try:
                    overlay.line.color.rgb = RGBColor(0, 0, 0)
                    overlay.line.width = 0
                except:
                    pass
            
            # 7. Move overlay to just above the background image
            try:
                overlay_elm = overlay.element
                shape_tree = slide.shapes._spTree
                shape_tree.remove(overlay_elm)
                shape_tree.insert(3, overlay_elm)
            except (AttributeError, IndexError):
                print("Warning: Could not reorder overlay shape")
        
        print("Background image and overlay added successfully")
        return True
        
    except FileNotFoundError:
        print(f"Image file not accessible: {image_path}")
        return False
    except Exception as e:
        print(f"Failed to add background image: {e}")
        import traceback
        traceback.print_exc()
        return False


def _draw_title_slide(slide, prs, slide_data, theme, topic):
    #_add_image_as_background(slide, prs, slide_data.get('image_path'), overlay_alpha=0.3)
    
    # Add a subtle gradient overlay
    overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    fill = overlay.fill
    fill.gradient()
    fill.gradient_stops[0].position = 0
    fill.gradient_stops[0].color.rgb = RGBColor(0, 0, 0)
    fill.gradient_stops[0].color.brightness = 0.4
    fill.gradient_stops[1].position = 1
    fill.gradient_stops[1].color.rgb = RGBColor(0, 0, 0)
    fill.gradient_stops[1].color.brightness = 0.2
    overlay.line.fill.background()
    
    # Title text (larger, more prominent)
    tx_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(14), Inches(3))
    p = tx_box.text_frame.paragraphs[0]
    p.text = topic
    p.font.name = theme['font']
    p.font.size = Pt(80)
    p.font.bold = True
    p.font.color.rgb = theme['text']
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle (minimal)
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(6), Inches(14), Inches(1))
    p = subtitle_box.text_frame.paragraphs[0]
    p.text = "An AI-Generated Presentation"
    p.font.name = theme['font']
    p.font.size = Pt(28)
    p.font.color.rgb = theme['text']
    p.alignment = PP_ALIGN.CENTER

def _draw_toc_slide(slide, prs, slides, theme):
    # Add a relevant background image
    #if slides and slides[0].get('image_path'):
        #_add_image_as_background(slide, prs, slides[0]['image_path'])
    
    # Add semi-transparent overlay for better readability
    overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = theme['bg']
    overlay.fill.fore_color.brightness = 0.8
    overlay.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(14), Inches(1))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = "Table of Contents"
    p.font.name = theme['font']
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = theme['text']
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(2), Inches(2.5), Inches(12), Inches(5))
    content_frame = content_box.text_frame
    
    for i, slide_data in enumerate(slides, 1):
        p = content_frame.add_paragraph()
        p.text = f"{i}. {slide_data['slide_title']}"
        p.font.name = theme['font']
        p.font.size = Pt(24)
        p.font.color.rgb = theme['text']
        p.space_after = Pt(12)

def _draw_photo_slide(slide, prs, slide_data, theme):
    _add_image_as_background(slide, prs, slide_data.get('image_path'))
    
    # Add a subtle gradient overlay for better text readability
    # overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    # fill = overlay.fill
    # fill.gradient()
    # fill.gradient_stops[0].position = 0
    # fill.gradient_stops[0].color.rgb = RGBColor(0, 0, 0)
    # fill.gradient_stops[0].color.brightness = 0.3
    # fill.gradient_stops[1].position = 1
    # fill.gradient_stops[1].color.rgb = RGBColor(0, 0, 0)
    # fill.gradient_stops[1].color.brightness = 0.1
    # overlay.line.fill.background()
    
    # Function to add text with markdown bold formatting
    def add_formatted_text(paragraph, text, font_size=Pt(32), is_title=False):
        # Split text by asterisks to find bold sections
        parts = text.split('*')
        
        for i, part in enumerate(parts):
            if not part:  # Skip empty parts
                continue
                
            if i == 0:
                # First part goes to the main paragraph
                paragraph.text = part
                paragraph.font.name = theme['font']
                paragraph.font.size = font_size
                paragraph.font.color.rgb = theme['text']
                if is_title:
                    paragraph.font.bold = True
                paragraph.alignment = PP_ALIGN.CENTER
            else:
                # Add subsequent parts as runs
                run = paragraph.add_run()
                run.text = part
                run.font.name = theme['font']
                run.font.size = font_size
                run.font.color.rgb = theme['text']
                
                # Make bold if it's an odd index (text between asterisks) or if it's a title
                if i % 2 == 1 or is_title:
                    run.font.bold = True
    
    # Title text (larger, more prominent)
    tx_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(14), Inches(2))
    tx_frame = tx_box.text_frame
    
    title_text = slide_data.get('slide_title', '')
    add_formatted_text(tx_frame.paragraphs[0], title_text, Pt(54), is_title=True)
    
    # Body text (single line, minimal)
    body_box = slide.shapes.add_textbox(Inches(1), Inches(7), Inches(14), Inches(1))
    body_frame = body_box.text_frame
    
    # Take only the first key point and make it a one-liner
    body_text = slide_data.get('slide_body', '')
    key_points = [point.strip() for point in body_text.split('.') if point.strip()]
    if key_points:
        add_formatted_text(body_frame.paragraphs[0], key_points[0], Pt(32))
    
    # Add supporting images if available
    supporting_images = slide_data.get('supporting_images', [])
    if supporting_images:
        print(f"DEBUG: Adding supporting images: {supporting_images}")
        
        # Calculate grid layout for supporting images
        num_images = len(supporting_images)
        if num_images == 1:
            # Single large image
            img_height = Inches(4.3)
            img_width = Inches(7)
            img_left = Inches(3.5)
            img_top = Inches(2.4)
            slide.shapes.add_picture(supporting_images[0], img_left, img_top, width=img_width, height=img_height)
        elif num_images == 2:
            # Two images side by side
            img_height = Inches(3.9)
            img_width = Inches(6)
            img_gap = Inches(1)
            img_top = Inches(2.4)
            
            # First image
            slide.shapes.add_picture(supporting_images[0], Inches(1.5), img_top, width=img_width, height=img_height)
            # Second image
            slide.shapes.add_picture(supporting_images[1], Inches(9.5), img_top, width=img_width, height=img_height)
        else:
            # Three images in a grid
            img_height = Inches(3.5)
            img_width = Inches(4.5)
            img_gap = Inches(0.5)
            img_top = Inches(2.4)
            
            for idx, img_path in enumerate(supporting_images[:3]):
                if not os.path.exists(img_path):
                    print(f"WARNING: Supporting image file not found: {img_path}")
                    continue
                col = idx % 3
                left = Inches(0.75 + (img_width + img_gap) * col)
                slide.shapes.add_picture(img_path, left, img_top, width=img_width, height=img_height)

    # Generate and add diagram
    diagram_type = _determine_diagram_type(slide_data)
    
    if os.environ.get("OPENAI_API_KEY"):
        diagram_path = _generate_diagram_with_openai(slide, slide_data, diagram_type)
    else:
        diagram_path = _generate_diagram(slide_data, diagram_type)

    if diagram_path and os.path.exists(diagram_path):
        # Add diagram in the remaining space
        diagram_width = Inches(4)
        diagram_height = Inches(3)
        diagram_left = Inches(10)
        diagram_top = Inches(2.5)
        slide.shapes.add_picture(diagram_path, diagram_left, diagram_top, width=diagram_width, height=diagram_height)

def _draw_diagram_slide(slide, prs, slide_data, theme):
    # Add a relevant background image
    if slide_data.get('background_image'):
        _add_image_as_background(slide, prs, slide_data['background_image'])
    else:
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = theme['bg']
    
    # Add decorative header
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.5))
    header.fill.solid()
    header.fill.fore_color.rgb = theme['accent']
    header.line.fill.background()
    
    # Title
    tx_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.25), Inches(7), Inches(1))
    tx_frame = tx_box.text_frame
    p = tx_frame.paragraphs[0]
    p.text = slide_data.get('slide_title', '')
    p.font.name = theme['font']
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = theme['text']
    
    # Content with automatic formatting and markdown bold handling
    content_box = slide.shapes.add_textbox(Inches(0.75), Inches(2), Inches(7), Inches(6))
    content_frame = content_box.text_frame
    content_text = slide_data.get('slide_body', '')
    
    # Function to add text with markdown bold formatting
    def add_formatted_text(paragraph, text):
        # Split text by asterisks to find bold sections
        parts = text.split('*')
        
        # Clear any existing text first
        paragraph.clear()
        
        for i, part in enumerate(parts):
            if not part:  # Skip empty parts
                continue
                
            if i == 0:
                # First part goes to the main paragraph
                paragraph.text = part
                paragraph.font.name = theme['font']
                paragraph.font.size = Pt(22)
                paragraph.font.color.rgb = theme['subtext']
                paragraph.line_spacing = 1.3
            else:
                # Add subsequent parts as runs
                run = paragraph.add_run()
                run.text = part
                run.font.name = theme['font']
                run.font.size = Pt(22)
                run.font.color.rgb = theme['subtext']
                
                # Make bold if it's an odd index (text between asterisks)
                if i % 2 == 1:
                    run.font.bold = True
    
    # Check if content is already formatted with numbers or bullets
    def is_already_formatted(text):
        lines = text.strip().split('\n')
        if len(lines) <= 1:
            return False
        
        # Check for numbered format (1., 2., etc.)
        numbered_pattern = re.compile(r'^\s*\d+\.\s+')
        # Check for bullet format (•, -, *, etc.)
        bullet_pattern = re.compile(r'^\s*[•\-\*]\s+')
        
        formatted_lines = 0
        for line in lines:
            if numbered_pattern.match(line) or bullet_pattern.match(line):
                formatted_lines += 1
        
        # If more than half the lines are formatted, consider it already formatted
        return formatted_lines > len(lines) / 2
    
    if is_already_formatted(content_text):
        # Content is already formatted, use as is but handle markdown bold
        add_formatted_text(content_frame.paragraphs[0], content_text)
    else:
        # Content needs formatting - split by periods and add numbering/bullets
        sentences = [s.strip() for s in content_text.split('.') if s.strip()]
        
        # Determine if we should use numbers or bullets
        # Use numbers if there are sequential steps, bullets for general points
        use_numbers = any(keyword in content_text.lower() for keyword in 
                         ['step', 'first', 'second', 'then', 'next', 'finally', 'process', 'method'])
        
        # Clear the default paragraph
        content_frame.clear()
        
        for i, sentence in enumerate(sentences):
            if i == 0:
                p = content_frame.paragraphs[0]
            else:
                p = content_frame.add_paragraph()
            
            formatted_sentence = f"{i + 1}. {sentence.strip()}" if use_numbers else f"• {sentence.strip()}"
            add_formatted_text(p, formatted_sentence)
            p.space_after = Pt(6)  # Add some space between points
    
    # Diagram
    if slide_data.get('image_path'):
        img_height = Inches(6)
        img_top = Inches(2)
        if not os.path.exists(slide_data['image_path']):
            print(f"WARNING: Diagram image file not found: {slide_data['image_path']}")
        else:
            slide.shapes.add_picture(slide_data['image_path'], Inches(8.25), img_top, height=img_height)

def _generate_diagram(slide_data: dict, diagram_type: str) -> Optional[str]:
    """Generate a diagram based on slide content and type."""
    try:
        # Ensure the output directory exists
        diagrams_dir = "output/diagrams"
        os.makedirs(diagrams_dir, exist_ok=True)
        
        # Create a unique filename based on content
        content_hash = hashlib.md5(json.dumps(slide_data).encode()).hexdigest()
        diagram_path = os.path.join(diagrams_dir, f"{content_hash}.png")
        
        # For now, use simple shapes to create diagrams
        if diagram_type == "flow":
            return _create_flow_diagram(slide_data, diagram_path)
        elif diagram_type == "comparison":
            return _create_comparison_diagram(slide_data, diagram_path)
        elif diagram_type == "timeline":
            return _create_timeline_diagram(slide_data, diagram_path)
        else:
            return _create_generic_diagram(slide_data, diagram_path)
    except Exception as e:
        logging.error(f"Error generating diagram: {str(e)}")
        return None


def _create_flow_diagram(slide_data: Dict, output_path: str) -> Optional[str]:
    """Create a flow diagram using matplotlib and save as PNG."""
    try:
        points = [p.strip() for p in slide_data.get('slide_body', '').split('.') if p.strip()]
        if not points:
            points = [slide_data.get('slide_title', 'Flow')]
        fig, ax = plt.subplots(figsize=(6, 2))
        ax.axis('off')
        n = len(points)
        for i, point in enumerate(points[:5]):
            x = 1 + i * 2
            y = 1
            box = plt.Rectangle((x, y), 1.5, 0.6, fc='#0070C0', ec='none', zorder=2)
            ax.add_patch(box)
            ax.text(x + 0.75, y + 0.3, point[:30] + ('...' if len(point) > 30 else ''), color='white', ha='center', va='center', fontsize=10, zorder=3)
            if i < n - 1:
                ax.arrow(x + 1.5, y + 0.3, 0.5, 0, head_width=0.15, head_length=0.2, fc='#0070C0', ec='#0070C0', zorder=1, length_includes_head=True)
        ax.set_xlim(0, 2 * max(3, n))
        ax.set_ylim(0, 2)
        plt.tight_layout()
        fig.savefig(output_path, bbox_inches='tight', transparent=True)
        plt.close(fig)
        return output_path
    except Exception as e:
        logging.error(f"Error creating flow diagram: {str(e)}")
        return None

def _create_comparison_diagram(slide_data: Dict, output_path: str) -> Optional[str]:
    """Create a comparison diagram using matplotlib and save as PNG."""
    try:
        points = [p.strip() for p in slide_data.get('slide_body', '').split('.') if p.strip()]
        if not points:
            points = [slide_data.get('slide_title', 'Comparison')]
        fig, ax = plt.subplots(figsize=(6, 2))
        ax.axis('off')
        for i, point in enumerate(points[:4]):
            x = 1 + (i % 2) * 3
            y = 1.5 - (i // 2) * 1
            box = plt.Rectangle((x, y), 2, 0.7, fc='#0070C0', ec='none', zorder=2)
            ax.add_patch(box)
            ax.text(x + 1, y + 0.35, point[:40] + ('...' if len(point) > 40 else ''), color='white', ha='center', va='center', fontsize=10, zorder=3)
        ax.set_xlim(0, 6)
        ax.set_ylim(0, 2.5)
        plt.tight_layout()
        fig.savefig(output_path, bbox_inches='tight', transparent=True)
        plt.close(fig)
        return output_path
    except Exception as e:
        logging.error(f"Error creating comparison diagram: {str(e)}")
        return None

def _create_timeline_diagram(slide_data: Dict, output_path: str) -> Optional[str]:
    """Create a timeline diagram using matplotlib and save as PNG."""
    try:
        points = [p.strip() for p in slide_data.get('slide_body', '').split('.') if p.strip()]
        if not points:
            points = [slide_data.get('slide_title', 'Timeline')]
        fig, ax = plt.subplots(figsize=(6, 2))
        ax.axis('off')
        n = len(points)
        ax.plot([1, 5], [1, 1], color='#0070C0', lw=3, zorder=1)
        for i, point in enumerate(points[:5]):
            x = 1 + (i + 1) * (4 / (n + 1))
            ax.plot(x, 1, 'o', color='#0070C0', markersize=12, zorder=2)
            ax.text(x, 1.2, point[:30] + ('...' if len(point) > 30 else ''), ha='center', va='bottom', fontsize=9, zorder=3)
        ax.set_xlim(0, 6)
        ax.set_ylim(0.5, 2)
        plt.tight_layout()
        fig.savefig(output_path, bbox_inches='tight', transparent=True)
        plt.close(fig)
        return output_path
    except Exception as e:
        logging.error(f"Error creating timeline diagram: {str(e)}")
        return None

def _create_generic_diagram(slide_data: Dict, output_path: str) -> Optional[str]:
    """Create a generic diagram using matplotlib and save as PNG."""
    try:
        points = [p.strip() for p in slide_data.get('slide_body', '').split('.') if p.strip()]
        if not points:
            points = [slide_data.get('slide_title', 'Concept')]
        fig, ax = plt.subplots(figsize=(4, 4))
        ax.axis('off')
        # Central node
        ax.add_patch(plt.Circle((2, 2), 0.5, color='#0070C0', zorder=2))
        ax.text(2, 2, slide_data.get('slide_title', '')[:20] + ('...' if len(slide_data.get('slide_title', '')) > 20 else ''), color='white', ha='center', va='center', fontsize=11, zorder=3)
        n = min(4, len(points))
        for i, point in enumerate(points[:4]):
            angle = i * (360 / n)
            x = 2 + 1.5 * math.cos(math.radians(angle))
            y = 2 + 1.5 * math.sin(math.radians(angle))
            ax.add_patch(plt.Rectangle((x - 0.7, y - 0.25), 1.4, 0.5, color='#0070C0', zorder=2))
            ax.text(x, y, point[:30] + ('...' if len(point) > 30 else ''), color='white', ha='center', va='center', fontsize=9, zorder=3)
            # Line from center to box
            ax.plot([2, x], [2, y], color='#0070C0', lw=2, zorder=1)
        ax.set_xlim(0, 4)
        ax.set_ylim(0, 4)
        plt.tight_layout()
        fig.savefig(output_path, bbox_inches='tight', transparent=True)
        plt.close(fig)
        return output_path
    except Exception as e:
        logging.error(f"Error creating generic diagram: {str(e)}")
        return None

def _generate_diagram_with_openai(slide, slide_data: dict, diagram_type: str) -> Optional[str]:
    """Generate a diagram using OpenAI's multimodal model (for future use)."""
    try:
        # Access the API key
        openai_api_key = os.environ.get("OPENAI_API_KEY")
        print("DEBUG: Attempting diagram generation with OpenAI GPT-4o.")
        # TODO: Implement OpenAI API call here using slide_data to generate diagram data or code
        return None # Return the path to the generated image file if successful
    except Exception as e:
        logging.error(f"Error generating diagram with OpenAI: {str(e)}")
        return None

def _determine_diagram_type(slide_data: dict) -> str:
    """Determine the most appropriate diagram type based on slide content."""
    title = slide_data.get('slide_title', '').lower()
    body = slide_data.get('slide_body', '').lower()
    
    if any(word in title or word in body for word in ['flow', 'process', 'steps', 'sequence']):
        return 'flow'
    elif any(word in title or word in body for word in ['compare', 'versus', 'vs', 'difference']):
        return 'comparison'
    elif any(word in title or word in body for word in ['timeline', 'history', 'future', 'roadmap']):
        return 'timeline'
    else:
        return 'generic'