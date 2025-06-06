from pptx import Presentation

def analyze_ppt_structure(pptx_path: str):
    """
    Opens a .pptx file and prints a structured report of its contents
    for debugging purposes.
    """
    print("\n--- PPT Analysis Report ---")
    try:
        prs = Presentation(pptx_path)
        for i, slide in enumerate(prs.slides):
            layout_name = slide.slide_layout.name
            print(f"\nSlide {i+1}: Layout '{layout_name}'")
            
            if not slide.shapes:
                print("  ...This slide has no shapes.")
                continue

            # Check for background image on the slide master
            if slide.slide_layout.background.fill.type == 6: # 6 is 'Picture'
                 print(f"  - Master Background: Contains an image.")

            # Check for shapes and placeholders on the slide itself
            for shape in slide.shapes:
                if shape.is_placeholder:
                    placeholder_type = shape.placeholder_format.type
                    shape_name = shape.name
                    if shape.has_text_frame and shape.text_frame.text:
                        print(f"  - Placeholder '{shape_name}' ({placeholder_type}): Contains text -> \"{shape.text_frame.text[:50]}...\"")
                    elif shape.shape_type == 13: # 13 is 'Picture'
                        print(f"  - Placeholder '{shape_name}' ({placeholder_type}): Contains an image.")
                    else:
                        print(f"  - Placeholder '{shape_name}' ({placeholder_type}): Is empty.")
                elif shape.shape_type == 13: # Non-placeholder picture
                     print(f"  - Shape (Picture): A non-placeholder image exists on this slide.")
    except Exception as e:
        print(f"Could not analyze the presentation. Error: {e}")