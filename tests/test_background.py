from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image
import tempfile

def create_test_slide():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank slide
    
    # Add background image
    img = Image.open('downloads/cache/Future of Renewable Energy solar energy landscape.png')
    if img.mode != 'RGB':
        img = img.convert('RGB')
    
    with tempfile.NamedTemporaryFile(suffix='.bmp', delete=False) as tmp_bmp:
        bmp_path = tmp_bmp.name
        img.save(bmp_path, 'BMP')
    
    # Add the image
    bg_pic = slide.shapes.add_picture(
        bmp_path, 0, 0, width=prs.slide_width, height=prs.slide_height
    )
    bg_pic.zorder = 0
    
    # Add semi-transparent overlay
    overlay = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    overlay.fill.transparency = 0.7  # 70% transparency
    overlay.zorder = 1
    
    # Save the presentation
    prs.save('test_background.pptx')
    print("Test presentation saved as test_background.pptx")

if __name__ == '__main__':
    create_test_slide() 