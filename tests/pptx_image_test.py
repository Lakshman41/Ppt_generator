from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from PIL import Image
import tempfile
import os

# Use one of the images from your cache
img_path = 'downloads/cache/Future of Renewable Energy wind power horizon.png'

# Convert to BMP for compatibility
with Image.open(img_path) as img:
    with tempfile.NamedTemporaryFile(suffix='.bmp', delete=False) as tmp_bmp:
        bmp_path = tmp_bmp.name
        img.convert('RGB').save(bmp_path, 'BMP')

prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)
blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)

# Add background image
slide.shapes.add_picture(bmp_path, 0, 0, width=prs.slide_width, height=prs.slide_height)

# Add a debug rectangle
rect = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.5), Inches(2), Inches(2)
)
rect.fill.transparency = 1.0  # Fully transparent
rect.line.color.rgb = RGBColor(255, 255, 255)

prs.save('test_output.pptx')
print('Minimal test presentation saved as test_output.pptx') 